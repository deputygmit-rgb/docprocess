import fitz
from pdf2image import convert_from_path
from pptx import Presentation
from openpyxl import load_workbook
from docx import Document as DocxDocument
from PIL import Image
import io
import base64
import os
from typing import List, Dict, Any
import re


def normalize_table(text: str) -> Dict[str, Any]:
    """
    Convert unstructured table text into normalized column structure with proper JSON rows.
    Handles various table formats and detects column structure automatically.
    """
    if not text or not text.strip():
        return {"columns": [], "rows": [], "row_count": 0, "column_count": 0}
    
    lines = [line.strip() for line in text.strip().split('\n') if line.strip()]
    
    if not lines:
        return {"columns": [], "rows": [], "row_count": 0, "column_count": 0}
    
    # Parse rows - try multiple delimiters
    parsed_rows = []
    for line in lines:
        # Try to split by common delimiters: | , tabs, or multiple spaces
        if '|' in line:
            cells = [c.strip() for c in line.split('|') if c.strip()]
        elif ',' in line and line.count(',') >= 2:
            cells = [c.strip() for c in line.split(',')]
        else:
            # Split by multiple spaces (2+) or tabs
            cells = re.split(r'\s{2,}|\t', line)
            cells = [c.strip() for c in cells if c.strip()]
        
        if cells:
            parsed_rows.append(cells)
    
    if not parsed_rows:
        return {"columns": [], "rows": [], "row_count": 0, "column_count": 0}
    
    # Determine the maximum number of columns
    max_cols = max(len(row) for row in parsed_rows) if parsed_rows else 1
    
    # Detect header row
    header = None
    data_start_idx = 0
    
    if len(parsed_rows) > 1:
        first_row = parsed_rows[0]
        second_row = parsed_rows[1]
        
        # Heuristics to detect if first row is a header:
        # 1. First row has mostly text (non-numeric), second row has numbers
        # 2. First row is shorter or padded with empty cells
        # 3. First row contains common header keywords
        
        first_numeric_count = sum(1 for cell in first_row if any(c.isdigit() for c in cell))
        second_numeric_count = sum(1 for cell in second_row if any(c.isdigit() for c in cell))
        
        header_keywords = {'name', 'id', 'date', 'value', 'amount', 'total', 'description', 'title', 
                          'category', 'type', 'status', 'count', 'item', 'col', 'column'}
        
        has_header_keyword = any(
            any(kw in cell.lower() for kw in header_keywords) 
            for cell in first_row
        )
        
        is_likely_header = (
            (first_numeric_count < second_numeric_count) or
            has_header_keyword or
            len(first_row) < len(second_row)
        )
        
        if is_likely_header:
            header = first_row
            data_start_idx = 1
    
    # Use generic headers if not detected
    if header is None:
        header = [f"Column_{i+1}" for i in range(max_cols)]
        data_start_idx = 0
    else:
        # Pad header to match column count
        header = (header + [''] * max_cols)[:max_cols]
    
    # Normalize data rows to match column count
    normalized_rows = []
    for row in parsed_rows[data_start_idx:]:
        # Pad or trim to match column count
        padded_row = (row + [''] * max_cols)[:max_cols]
        
        # Create dictionary mapping headers to values
        row_dict = {}
        for i, col_header in enumerate(header):
            col_name = col_header if col_header else f"Column_{i+1}"
            row_dict[col_name] = padded_row[i]
        
        normalized_rows.append(row_dict)
    
    return {
        "columns": header,
        "rows": normalized_rows,
        "row_count": len(normalized_rows),
        "column_count": len(header),
        "detection_info": {
            "max_columns_detected": max_cols,
            "header_detected": data_start_idx > 0,
            "parser_type": "column_normalized"
        },
        "raw_array_format": {
            "columns": header,
            "data": [
                [row.get(col, '') for col in header]
                for row in normalized_rows
            ]
        }
    }


class DocumentProcessor:
    
    @staticmethod
    def pdf_to_images(pdf_path: str) -> List[Image.Image]:
        try:
            images = convert_from_path(pdf_path, dpi=150)
            return images
        except Exception as e:
            raise Exception(f"PDF conversion failed: {str(e)}")
    
    @staticmethod
    def pdf_extract_pages(pdf_path: str) -> List[Dict[str, Any]]:
        pages = []
        doc = fitz.open(pdf_path)
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            img_data = pix.tobytes("png")
            img = Image.open(io.BytesIO(img_data))
            
            text = page.get_text()
            
            pages.append({
                "page_number": page_num + 1,
                "image": img,
                "text": text,
                "width": page.rect.width,
                "height": page.rect.height
            })
        
        doc.close()
        return pages
    
    @staticmethod
    def ppt_to_images(ppt_path: str) -> List[Dict[str, Any]]:
        prs = Presentation(ppt_path)
        slides = []
        
        for idx, slide in enumerate(prs.slides):
            text_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content.append(shape.text)
            
            slides.append({
                "slide_number": idx + 1,
                "text": "\n".join(text_content),
                "shapes_count": len(slide.shapes)
            })
        
        return slides
    
    @staticmethod
    def xlsx_extract_data(xlsx_path: str) -> List[Dict[str, Any]]:
        wb = load_workbook(xlsx_path, data_only=True)
        sheets_data = []
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            
            data = []
            for row in ws.iter_rows(values_only=True):
                if any(cell is not None for cell in row):
                    data.append(list(row))
            
            sheets_data.append({
                "sheet_name": sheet_name,
                "data": data,
                "rows": len(data),
                "columns": ws.max_column
            })
        
        wb.close()
        return sheets_data
    
    @staticmethod
    def docx_extract_text(docx_path: str) -> Dict[str, Any]:
        doc = DocxDocument(docx_path)
        
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        
        tables_data = []
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            tables_data.append(table_data)
        
        return {
            "paragraphs": paragraphs,
            "tables": tables_data,
            "paragraph_count": len(paragraphs),
            "table_count": len(tables_data)
        }
    
    @staticmethod
    def image_to_base64(image: Image.Image) -> str:
        buffered = io.BytesIO()
        image.save(buffered, format="PNG")
        img_str = base64.b64encode(buffered.getvalue()).decode()
        return f"data:image/png;base64,{img_str}"
    
    @staticmethod
    def image_to_base64_from_file(file_path: str) -> str:
        """Convert an image file (JPG, PNG, etc.) directly to base64 data URI"""
        try:
            image = Image.open(file_path)
            # Determine the image format
            image_format = image.format or "PNG"
            if image_format.upper() in ['JPG', 'JPEG']:
                image_format = "JPEG"
            
            buffered = io.BytesIO()
            image.save(buffered, format=image_format)
            img_str = base64.b64encode(buffered.getvalue()).decode()
            mime_type = f"image/{image_format.lower()}"
            if mime_type == "image/jpeg":
                mime_type = "image/jpeg"
            return f"data:{mime_type};base64,{img_str}"
        except Exception as e:
            raise ValueError(f"Failed to process image file {file_path}: {str(e)}")